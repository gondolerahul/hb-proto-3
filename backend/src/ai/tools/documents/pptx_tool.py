"""
PPTX Tool — Read, Create, and Update PowerPoint presentations for AI Agents.

Uses python-pptx to manipulate .pptx files. Saves generated files
via ArtifactService for traceability.
"""

import json
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional

from src.ai.tools.base import Tool

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None

BASE_DIR = Path(__file__).resolve().parents[3] / "artifact" / "system-generated"


class PptxTool(Tool):
    """Tool for reading, creating, and updating PowerPoint (.pptx) presentations.

    Actions:
        create  — Build a new pptx with title slide, content slides, bullet slides
        read    — Extract slide text, notes, and metadata from an existing pptx
        update  — Add new slides or modify text on existing slides
    """

    name = "pptx_tool"
    description = (
        "Read, create, or update PowerPoint (.pptx) presentations. "
        "Input: JSON with 'action' ('create', 'read', or 'update') "
        "and action-specific parameters. "
        "For 'create': provide 'filename', 'title', and 'slides' array. "
        "For 'read': provide 'file_path'. "
        "For 'update': provide 'file_path' and slides to add/modify."
    )

    async def run(self, input_data: str) -> str:
        return await self.run_with_context(input_data, context=None)

    async def run_with_context(self, input_data: str, context=None) -> str:
        if Presentation is None:
            return json.dumps({"error": "python-pptx is not installed"})

        try:
            params = json.loads(input_data) if isinstance(input_data, str) else input_data
            action = params.get("action")

            if action == "create":
                return await self._create(params, context)
            elif action == "read":
                return self._read(params)
            elif action == "update":
                return await self._update(params, context)
            else:
                return json.dumps({"error": f"Unknown action: {action}. Use 'create', 'read', or 'update'."})
        except Exception as e:
            return json.dumps({"error": f"PptxTool error: {str(e)}"})

    # ------------------------------------------------------------------
    # CREATE
    # ------------------------------------------------------------------

    async def _create(self, params: Dict[str, Any], context: Optional[Dict[str, Any]]) -> str:
        prs = Presentation()

        title = params.get("title", "Untitled Presentation")
        subtitle = params.get("subtitle", "")
        filename = params.get("filename", "presentation.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        # Title slide
        title_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = subtitle

        # Content slides
        slides_data = params.get("slides", [])
        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")

            if slide_type == "title_only":
                layout = prs.slide_layouts[5]  # Blank
                s = prs.slides.add_slide(layout)
                if slide_data.get("title"):
                    txBox = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = slide_data["title"]
                    p.font.size = Pt(28)
                    p.font.bold = True

            elif slide_type == "bullets":
                layout = prs.slide_layouts[1]  # Title and Content
                s = prs.slides.add_slide(layout)
                if slide_data.get("title"):
                    s.shapes.title.text = slide_data["title"]
                body = s.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(slide_data.get("bullets", [])):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                    tf.paragraphs[i].level = slide_data.get("bullet_level", 0)

            elif slide_type == "two_column":
                layout = prs.slide_layouts[3]  # Two Content
                s = prs.slides.add_slide(layout)
                if slide_data.get("title"):
                    s.shapes.title.text = slide_data["title"]
                left_content = slide_data.get("left_content", "")
                right_content = slide_data.get("right_content", "")
                if len(s.placeholders) > 1:
                    s.placeholders[1].text = left_content
                if len(s.placeholders) > 2:
                    s.placeholders[2].text = right_content

            else:  # content (default)
                layout = prs.slide_layouts[1]  # Title and Content
                s = prs.slides.add_slide(layout)
                if slide_data.get("title"):
                    s.shapes.title.text = slide_data["title"]
                content = slide_data.get("content", "")
                if content:
                    s.placeholders[1].text = content

            # Add notes
            notes = slide_data.get("notes", "")
            if notes:
                s.notes_slide.notes_text_frame.text = notes

        # Save
        company_id = (context or {}).get("company_id", params.get("company_id", "default"))
        date_str = datetime.utcnow().strftime("%Y-%m-%d")
        output_dir = BASE_DIR / str(company_id) / date_str
        output_dir.mkdir(parents=True, exist_ok=True)

        file_path = output_dir / filename
        prs.save(str(file_path))

        artifact_id = await self._register_artifact(file_path, filename, company_id, params)

        result = {
            "status": "success",
            "file_path": str(file_path),
            "message": f"Created {filename} with {len(slides_data) + 1} slides",
        }
        if artifact_id:
            result["artifact_id"] = artifact_id
            result["download_url"] = f"/api/v1/artifacts/{artifact_id}/download"
            result["document_path"] = f"/api/v1/artifacts/{artifact_id}/download"

        return json.dumps(result)

    # ------------------------------------------------------------------
    # READ
    # ------------------------------------------------------------------

    def _read(self, params: Dict[str, Any]) -> str:
        file_path = params.get("file_path", "")
        if not file_path:
            return json.dumps({"error": "Missing 'file_path'"})

        if not os.path.exists(file_path):
            fallback = os.path.join("uploads", os.path.basename(file_path))
            if os.path.exists(fallback):
                file_path = fallback
            else:
                return json.dumps({"error": f"File not found: {file_path}"})

        prs = Presentation(file_path)

        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_info = {"slide_number": idx + 1, "shapes": []}

            for shape in slide.shapes:
                shape_data = {"shape_type": shape.shape_type, "name": shape.name}
                if shape.has_text_frame:
                    shape_data["text"] = shape.text_frame.text
                if hasattr(shape, "table"):
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text for cell in row.cells])
                    shape_data["table"] = rows
                slide_info["shapes"].append(shape_data)

            # Notes
            try:
                if slide.has_notes_slide:
                    slide_info["notes"] = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass

            slides.append(slide_info)

        return json.dumps({
            "status": "success",
            "total_slides": len(slides),
            "slides": slides,
            "slide_width": str(prs.slide_width),
            "slide_height": str(prs.slide_height),
        })

    # ------------------------------------------------------------------
    # UPDATE
    # ------------------------------------------------------------------

    async def _update(self, params: Dict[str, Any], context: Optional[Dict[str, Any]]) -> str:
        file_path = params.get("file_path", "")
        if not file_path or not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"})

        prs = Presentation(file_path)
        changes_made = 0

        # Add new slides
        new_slides = params.get("add_slides", [])
        for slide_data in new_slides:
            slide_type = slide_data.get("type", "content")

            if slide_type == "bullets":
                layout = prs.slide_layouts[1]
                s = prs.slides.add_slide(layout)
                if slide_data.get("title"):
                    s.shapes.title.text = slide_data["title"]
                body = s.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(slide_data.get("bullets", [])):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
            else:
                layout = prs.slide_layouts[1]
                s = prs.slides.add_slide(layout)
                if slide_data.get("title"):
                    s.shapes.title.text = slide_data["title"]
                content = slide_data.get("content", "")
                if content:
                    s.placeholders[1].text = content

            notes = slide_data.get("notes", "")
            if notes:
                s.notes_slide.notes_text_frame.text = notes
            changes_made += 1

        # Modify existing slide text
        modifications = params.get("modify_slides", [])
        for mod in modifications:
            slide_idx = mod.get("slide_number", 0) - 1
            if 0 <= slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                new_title = mod.get("title")
                if new_title and slide.shapes.title:
                    slide.shapes.title.text = new_title
                    changes_made += 1
                new_content = mod.get("content")
                if new_content:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            shape.text_frame.text = new_content
                            changes_made += 1
                            break

        save_path = params.get("save_as", file_path)
        prs.save(save_path)

        return json.dumps({
            "status": "success",
            "file_path": save_path,
            "changes_made": changes_made,
            "total_slides": len(prs.slides),
            "message": f"Updated presentation with {changes_made} changes",
        })

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    async def _register_artifact(
        file_path: Path, filename: str, company_id: str, params: Dict[str, Any]
    ) -> str | None:
        """Register file as an artifact. Returns artifact ID or None."""
        if company_id and company_id != "default":
            try:
                from uuid import UUID as _UUID
                from src.common.database import AsyncSessionLocal
                from src.ai.artifact_service import ArtifactService, ORIGIN_SYSTEM

                async with AsyncSessionLocal() as _db:
                    art_svc = ArtifactService(_db)
                    with open(file_path, "rb") as f:
                        file_bytes = f.read()
                    artifact = await art_svc.save_artifact(
                        file_bytes=file_bytes,
                        file_name=filename,
                        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        file_category="documents",
                        origin=ORIGIN_SYSTEM,
                        company_id=_UUID(str(company_id)),
                        purpose=params.get("purpose", "AI-generated PowerPoint presentation"),
                        generated_by=params.get("generated_by", "pptx_tool"),
                    )
                    return str(artifact.id)
            except Exception:
                pass  # Non-fatal
        return None

    def get_function_schema(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "enum": ["create", "read", "update"],
                        "description": "Action to perform",
                    },
                    "filename": {
                        "type": "string",
                        "description": "Output filename for 'create' (e.g. deck.pptx)",
                    },
                    "title": {
                        "type": "string",
                        "description": "Presentation title for 'create'",
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Title slide subtitle for 'create'",
                    },
                    "slides": {
                        "type": "array",
                        "description": (
                            "Array of slide objects for 'create'. Each slide: "
                            "'type' ('content'|'bullets'|'title_only'|'two_column'), "
                            "'title', 'content', 'bullets' (str[]), 'notes'"
                        ),
                        "items": {"type": "object"},
                    },
                    "file_path": {
                        "type": "string",
                        "description": "Path to existing pptx for 'read' or 'update'",
                    },
                    "add_slides": {
                        "type": "array",
                        "description": "Slides to add (for 'update')",
                        "items": {"type": "object"},
                    },
                    "modify_slides": {
                        "type": "array",
                        "description": (
                            "Modifications to existing slides (for 'update'): "
                            "[{slide_number, title, content}]"
                        ),
                        "items": {"type": "object"},
                    },
                    "save_as": {
                        "type": "string",
                        "description": "Optional save-as path (for 'update')",
                    },
                },
                "required": ["action"],
            },
        }
