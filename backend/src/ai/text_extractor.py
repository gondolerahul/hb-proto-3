"""
text_extractor.py — Unified text extraction from binary file formats.

Extracted from ExecutionEngine._extract_text_from_file() during Phase 6
deduplication to enable reuse across worker.py and service.py without
coupling to the engine instance.
"""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text_from_file(file_path, mime_type: str = "") -> str:
    """
    Extract human-readable text from a file.

    Handles binary document formats (DOCX, PDF, XLSX, PPTX) via their
    respective libraries, and falls back to plain text reading for
    everything else. Returns sanitized text with null bytes stripped.

    Args:
        file_path: Path to the file (str or Path).
        mime_type: Optional MIME type hint.

    Returns:
        Extracted text content, sanitized for PostgreSQL.
    """
    _fpath = Path(file_path)
    _ext = _fpath.suffix.lower()

    try:
        if _ext == ".docx":
            import docx
            doc = docx.Document(str(_fpath))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif _ext == ".pdf":
            import PyPDF2
            with open(_fpath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = "\n".join(
                    page.extract_text() or "" for page in reader.pages
                )
        elif _ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(_fpath), read_only=True, data_only=True)
                rows = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    rows.append(f"--- Sheet: {sheet} ---")
                    for row in ws.iter_rows(max_row=500, values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        rows.append(" | ".join(cells))
                wb.close()
                text = "\n".join(rows)
            except ImportError:
                text = _fpath.read_text(errors="replace")
        elif _ext in (".csv", ".tsv"):
            text = _fpath.read_text(errors="replace")
        elif _ext in (".pptx",):
            try:
                from pptx import Presentation
                prs = Presentation(str(_fpath))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_parts = [f"--- Slide {i} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_parts.append(shape.text)
                    slides_text.append("\n".join(slide_parts))
                text = "\n\n".join(slides_text)
            except ImportError:
                text = _fpath.read_text(errors="replace")
        else:
            # Plain text files (txt, md, json, xml, html, etc.)
            text = _fpath.read_text(errors="replace")
    except Exception as e:
        logger.warning(f"Text extraction failed for {_fpath.name}: {e}")
        try:
            text = _fpath.read_text(errors="replace")
        except Exception:
            return ""

    # Sanitize: remove null bytes that PostgreSQL UTF-8 encoding rejects
    text = text.replace("\x00", "")
    return text.strip()
